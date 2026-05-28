"""Genera SMECO_campus_robotizzato.pptx: slide con testo, foto e video (layout uniforme)."""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

_REPO = Path(__file__).resolve().parent.parent
OUT_REPO = _REPO / "SMECO_campus_robotizzato.pptx"
OUT_REPO_ALT = _REPO / "SMECO_campus_robotizzato_nuovo.pptx"
OUT_DOWNLOADS = Path.home() / "Downloads" / "SMECO_campus_robotizzato.pptx"

# Tipografia da relazione (Office/Windows)
FONT_TITLE = "Cambria"
FONT_BODY = "Calibri"

NAVY = RGBColor(0x1A, 0x2B, 0x4C)
GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT = RGBColor(0x00, 0x33, 0x66)

GROUP = "Gruppo SMECO"
TEAM_LIST = (
    "Luca Furlan, Giovanni Pitton, Leonardo Furlan, "
    "Mattia Scevola, Marco Iannone, Filippo Mariuzzo"
)


def _font(run, pt, bold=False, italic=False, color=None, name=None):
    run.font.size = Pt(pt)
    run.font.bold = bold
    run.font.italic = italic
    if name:
        run.font.name = name
    if color is not None:
        run.font.color.rgb = color


def _rule(slide, top_inches):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(top_inches), Inches(9.1), Inches(0.035)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()


def add_slide_three_blocks(
    prs,
    title: str,
    presentation_paragraphs: list[str],
    foto_spec: str,
    video_spec: str,
    index: int,
    total: int,
):
    """Ogni slide: titolo + Testo di presentazione + Foto da inserire + Video da inserire."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Titolo slide
    tb = slide.shapes.add_textbox(Inches(0.45), Inches(0.28), Inches(9.1), Inches(0.55))
    tf = tb.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    _font(r0, 20, bold=True, color=NAVY, name=FONT_TITLE)

    _rule(slide, 0.88)

    y = 1.02
    # --- Testo di presentazione ---
    h_text = 3.0
    box_t = slide.shapes.add_textbox(Inches(0.45), Inches(y), Inches(9.1), Inches(h_text))
    tft = box_t.text_frame
    tft.word_wrap = True
    p = tft.paragraphs[0]
    rl = p.add_run()
    rl.text = "Testo di presentazione"
    _font(rl, 11, bold=True, color=NAVY, name=FONT_TITLE)
    for para in presentation_paragraphs:
        p = tft.add_paragraph()
        p.space_before = Pt(5)
        rp = p.add_run()
        rp.text = para
        _font(rp, 12, color=GRAY, name=FONT_BODY)

    y += h_text + 0.08
    # --- Foto ---
    h_f = 0.92
    box_f = slide.shapes.add_textbox(Inches(0.45), Inches(y), Inches(9.1), Inches(h_f))
    tff = box_f.text_frame
    tff.word_wrap = True
    p = tff.paragraphs[0]
    rf = p.add_run()
    rf.text = "Foto da inserire"
    _font(rf, 11, bold=True, color=NAVY, name=FONT_TITLE)
    p = tff.add_paragraph()
    p.space_before = Pt(4)
    rf2 = p.add_run()
    rf2.text = foto_spec
    _font(rf2, 11, color=GRAY, name=FONT_BODY)

    y += h_f + 0.08
    # --- Video ---
    h_v = 0.92
    box_v = slide.shapes.add_textbox(Inches(0.45), Inches(y), Inches(9.1), Inches(h_v))
    tfv = box_v.text_frame
    tfv.word_wrap = True
    p = tfv.paragraphs[0]
    rv = p.add_run()
    rv.text = "Video da inserire"
    _font(rv, 11, bold=True, color=NAVY, name=FONT_TITLE)
    p = tfv.add_paragraph()
    p.space_before = Pt(4)
    rv2 = p.add_run()
    rv2.text = video_spec
    _font(rv2, 11, color=GRAY, name=FONT_BODY)

    foot = slide.shapes.add_textbox(Inches(0.45), Inches(7.02), Inches(9.1), Inches(0.38))
    tffo = foot.text_frame
    p = tffo.paragraphs[0]
    rf = p.add_run()
    rf.text = f"{GROUP} · Campus robotizzato · {index}/{total}"
    _font(rf, 9, color=RGBColor(0x66, 0x66, 0x66), name=FONT_BODY)

    return slide


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slides_data: list[tuple[str, list[str], str, str]] = [
        (
            f"{GROUP} — Campus robotizzato",
            [
                (
                    "Buongiorno. Siamo il Gruppo SMECO. La proposta integra quattro riferimenti di prodotto "
                    "Unitree: Go2 con braccio Z1 per manipolazione e terreni articolati; A2-W per mobilità su "
                    "ruote con adeguata protezione ambientale in esterno; G1-D per accoglienza e interazione; "
                    "B2-W per pattugliamento e ispezioni programmate."
                ),
                (
                    "Go2 e A2-W sono complementari sul campus: il quadrupede per accessi e compiti specifici, "
                    "la piattaforma A2-W dove serve tenuta IP in presenza di pioggia, polvere o servizio prolungato all’aperto."
                ),
                (
                    "L’impostazione resta graduale: pilota su una porzione del sito, con sicurezza, privacy "
                    "e ripetibilità delle missioni, prima di estensioni di flotta."
                ),
            ],
            "Logo del gruppo SMECO oppure fotografia di squadra con Go2, A2-W, G1-D e B2-W in ambiente reale.",
            "Breve sequenza di apertura (circa 5–10 secondi), oppure solo materiale fotografico.",
        ),
        (
            "Go2 e A2-W — ruoli sul campus",
            [
                (
                    "Il quadrupede Unitree Go2 resta il riferimento per aree verdi, cordoli e fondi misti, "
                    "dove serve locomozione articolata e, con braccio Z1, manipolazione leggera in zone assegnate."
                ),
                (
                    "Il Go2, nella configurazione di progetto, non dispone di un grado di protezione IP "
                    "adequato a un uso continuativo all’esposizione diretta a pioggia, spruzzi o polveri "
                    "intensive: non lo proponiamo come unica piattaforma per tutte le missioni esterne senza riparo."
                ),
                (
                    "Affianchiamo quindi la Unitree A2-W: piattaforma su ruote pensata per percorsi campus "
                    "con requisito di tenuta ambientale (IP), logistica tra edifici e servizio meteo senza "
                    "compromettere disponibilità operativa."
                ),
                (
                    "L’operatore definisce aree e missioni; a fine turno si documentano esiti, tempi e anomalie "
                    "in modo sintetico, separando dove competono Go2 e dove A2-W."
                ),
            ],
            "Composizione a due pannelli: Go2 su terreno misto; A2-W su viale o piazzale bagnato o in condizioni di luce pioggia (se disponibile).",
            "Due clip: Go2 su tratto impegnativo; A2-W su percorso esterno lungo il campus (15–25 secondi ciascuna o montaggio unico).",
        ),
        (
            "Go2 e braccio Z1 — manipolazione",
            [
                (
                    "Il braccio Z1, con pinza parallela o ventosa, integra la funzione manipolativa: "
                    "dalla stima della posa dell’oggetto alla pianificazione e all’esecuzione della presa."
                ),
                (
                    "In contesto campus intervengono vento, limiti di carico, geometrie non catalogate "
                    "e la presenza di pedoni. La roadmap prevede validazione in simulazione, quindi prove "
                    "in area ristretta e controllata."
                ),
            ],
            "Dettaglio di Z1 e end-effector in laboratorio o in campo durante una prova autorizzata.",
            "Sequenza breve di avvicinamento controllato all’oggetto, se disponibile e conforme alle norme di sicurezza.",
        ),
        (
            "Simulazione MuJoCo — confronto tra regimi di presa",
            [
                (
                    "Nel simulatore MuJoCo abbiamo riprodotto il compito di presa su un oggetto semplice, "
                    "come passaggio obbligato prima delle prove sul robot reale."
                ),
                (
                    "Mostriamo due regimi: presa con forza contenuta e comportamento più compliant, "
                    "e presa più decisa. Stesso scenario virtuale, per discutere stabilità, rischio di "
                    "danneggiamento dell’oggetto e ripetibilità."
                ),
            ],
            "Due immagini affiancate: screenshot MuJoCo per presa leggera e per presa più forte.",
            "Due clip brevi oppure un’unica clip con confronto tra le due configurazioni di presa.",
        ),
        (
            "Accoglienza e spostamento — G1-D e A2-W",
            [
                (
                    "Per l’accoglienza e l’interazione ravvicinata utilizziamo esplicitamente la variante "
                    "Unitree G1-D: umanoide come punto di contatto mobile per percorsi, orari e servizi, "
                    "con risposta vocale e dominio informativo delimitato."
                ),
                (
                    "Per gli spostamenti orizzontali efficienti e per le condizioni esterne che richiedono "
                    "protezione IP, la flotta campus include la Unitree A2-W su vie e cortili, in coerenza "
                    "con quanto illustrato sul binomio Go2–A2-W."
                ),
                (
                    "L’accoppiamento G1-D e A2-W separa nettamente interazione umana di qualità da logistica "
                    "su ruote protetta, riducendo sovrapposizioni di ruolo sul campo."
                ),
            ],
            "G1-D in atrio o area coperta o esterna controllata; A2-W su percorso pavimentato; composizione a due pannelli affiancati.",
            "G1-D durante benvenuto o risposta breve; A2-W in transizione tra due punti del campus (circa 10–20 secondi ciascuno o montaggio unico).",
        ),
        (
            "Interazione vocale — G1-D",
            [
                (
                    "Sulla piattaforma G1-D illustriamo il flusso completo: voce in ingresso, riconoscimento "
                    "del parlato, ragionamento tramite modello linguistico con eventuale recupero da documenti "
                    "aggiornati, sintesi vocale in uscita."
                ),
                (
                    "Le registrazioni proposte sono state realizzate da Lorenzon: mostrano domande plausibili "
                    "poste dal gruppo e le risposte del sistema in una condizione operativa rappresentativa sul G1-D."
                ),
            ],
            "G1-D in ascolto o schermata dell’interfaccia; in alternativa schema essenziale del flusso voce–elaborazione–voce.",
            "Registrazioni integrali della dimostrazione vocale sul G1-D: file principale di questa slide.",
        ),
        (
            "Pattugliamento — Unitree B2-W",
            [
                (
                    "Il pattugliamento è assegnato alla piattaforma su ruote B2-W, con missioni su waypoints "
                    "e percezione a lungo raggio tramite LiDAR."
                ),
                (
                    "Telecamere complementari, ove autorizzate, possono alimentare log e verifiche successive, "
                    "nel rispetto della normativa privacy e delle policy del sito. Obiettivo: tracciabilità dei giri "
                    "e segnalazione di situazioni anomale lungo il percorso."
                ),
            ],
            "B2-W in corsia o viale; inquadratura che evidenzi il profilo del veicolo o il sensore principale.",
            "Pattugliamento su un tratto reale del campus; durata indicativa 20–40 secondi.",
        ),
        (
            "G1-D — apprendimento, istruzioni e realtà virtuale",
            [
                (
                    "Sull’umanoide G1-D esploriamo reinforcement learning e comportamenti guidati da istruzioni "
                    "in linguaggio naturale, con vincoli di sicurezza e verifiche ripetibili sulle varianti hardware "
                    "e software della linea D."
                ),
                (
                    "La realtà virtuale supporta teleoperazione e raccolta dati in scenario controllato; "
                    "il materiale fotografico documenta le sessioni. Dove disponibile, affianchiamo estratti "
                    "da simulazione umanoide al materiale reale sul G1-D."
                ),
                (
                    "Attività immersive e documentazione correlate: contributo di Samuele su teleoperazione, "
                    "gestione in VR e riprese statiche di supporto."
                ),
            ],
            "Operatore in sessione VR; oppure G1-D in laboratorio; oppure frame da simulatore umanoide.",
            "Estratto di teleoperazione in VR o clip breve da simulazione; se non disponibile, solo materiale fotografico.",
        ),
        (
            "Grazie",
            [
                (
                    "Grazie per l’attenzione. Il Gruppo SMECO è disponibile per domande su perimetro di pilota, "
                    "sicurezza, integrazione con i servizi esistenti e roadmap di scala."
                ),
                (f"Componenti: {TEAM_LIST}."),
            ],
            "Fotografia di gruppo SMECO sul campo o vista rappresentativa del campus.",
            "Nessun video; passaggio diretto al dibattito con il pubblico.",
        ),
    ]

    total = len(slides_data)
    for i, (title, paras, foto, video) in enumerate(slides_data, start=1):
        add_slide_three_blocks(prs, title, paras, foto, video, i, total)

    try:
        prs.save(OUT_REPO)
        out_main = OUT_REPO
        print("OK — PPTX progetto:", OUT_REPO.resolve())
    except PermissionError:
        prs.save(OUT_REPO_ALT)
        out_main = OUT_REPO_ALT
        print(
            "File principale bloccato (chiudi SMECO_campus_robotizzato.pptx). Scritto:",
            OUT_REPO_ALT.resolve(),
        )
    try:
        prs.save(OUT_DOWNLOADS)
        print("OK — PPTX Download:", OUT_DOWNLOADS.resolve())
    except PermissionError:
        print("Download non scrivibile.")

    return out_main


if __name__ == "__main__":
    main()
